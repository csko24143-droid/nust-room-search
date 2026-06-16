from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Color palette
C_TITLE_BG   = RGBColor(0x1A, 0x53, 0x76)   # deep blue
C_PART1_BG   = RGBColor(0xFF, 0x8C, 0x00)   # orange  (糸電話)
C_PART2_BG   = RGBColor(0x1E, 0x8B, 0xC8)   # sky blue (気柱共鳴)
C_PART3_BG   = RGBColor(0x27, 0xAE, 0x60)   # green   (カリンバ)
C_SUMMARY_BG = RGBColor(0x8E, 0x44, 0xAD)   # purple  (まとめ)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK       = RGBColor(0x22, 0x22, 0x22)
C_LIGHT_BG   = RGBColor(0xF4, 0xF7, 0xFB)
C_ACCENT     = RGBColor(0xFF, 0xD7, 0x00)   # yellow accent

blank_layout = prs.slide_layouts[6]  # completely blank

def add_rect(slide, left, top, width, height, color, transparency=0):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, text, left, top, width, height,
                 font_size=24, bold=False, color=C_DARK,
                 align=PP_ALIGN.LEFT, wrap=True, font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_header_bar(slide, bg_color, part_label, title_text, icon=""):
    add_rect(slide, 0, 0, 13.33, 1.5, bg_color)
    if part_label:
        add_text_box(slide, part_label, 0.3, 0.05, 4, 0.45,
                     font_size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.LEFT)
    full_title = icon + " " + title_text if icon else title_text
    add_text_box(slide, full_title, 0.3, 0.45, 12.5, 0.9,
                 font_size=30, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

def add_slide(prs, bg_color, part_label, title, icon, content_fn):
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 7.5, C_LIGHT_BG)
    add_header_bar(slide, bg_color, part_label, title, icon)
    content_fn(slide)
    return slide

def bullet_box(slide, items, left=0.5, top=1.7, width=12.3, height=5.5,
               font_size=22, color=C_DARK, circle_color=None, prefix="●"):
    y = top
    line_h = height / len(items)
    for item in items:
        bullet = prefix + "  " if prefix else ""
        add_text_box(slide, bullet + item, left, y, width, line_h,
                     font_size=font_size, color=color)
        y += line_h

def two_col(slide, left_items, right_items,
            left_head="", right_head="",
            top=1.7, col_w=5.9, font_size=20):
    if left_head:
        add_text_box(slide, left_head, 0.4, top, col_w, 0.5,
                     font_size=20, bold=True, color=C_DARK)
    if right_head:
        add_text_box(slide, right_head, 6.9, top, col_w, 0.5,
                     font_size=20, bold=True, color=C_DARK)
    y = top + (0.55 if left_head or right_head else 0)
    lh = (5.5 - (y - top)) / max(len(left_items), len(right_items))
    for i, item in enumerate(left_items):
        add_text_box(slide, "● " + item, 0.5, y + i*lh, col_w, lh, font_size=font_size)
    for i, item in enumerate(right_items):
        add_text_box(slide, "● " + item, 7.0, y + i*lh, col_w, lh, font_size=font_size)
    # divider
    div = slide.shapes.add_shape(1, Inches(6.6), Inches(top), Inches(0.05), Inches(5.5))
    div.fill.solid(); div.fill.fore_color.rgb = RGBColor(0xCC,0xCC,0xCC)
    div.line.fill.background()

def step_box(slide, steps, top=1.7, width=11.8, left=0.8, font_size=21):
    h = 5.2 / len(steps)
    for i, step in enumerate(steps):
        y = top + i * h
        add_rect(slide, left - 0.4, y + 0.05, 0.55, h - 0.15,
                 RGBColor(0x1A,0x53,0x76))
        add_text_box(slide, str(i+1), left - 0.32, y + 0.08, 0.4, h - 0.15,
                     font_size=font_size-2, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, step, left + 0.2, y + 0.05, width, h - 0.15,
                     font_size=font_size, color=C_DARK)

def highlight_box(slide, text, left=0.5, top=2.5, width=12.3, height=1.5,
                  bg=RGBColor(0xFF,0xF3,0xCD), font_size=28):
    add_rect(slide, left, top, width, height, bg)
    add_text_box(slide, text, left + 0.2, top + 0.1, width - 0.4, height - 0.2,
                 font_size=font_size, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

# ─── SLIDE 1: タイトル ───────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, C_TITLE_BG)
# wave decoration
add_rect(slide, 0, 5.5, 13.33, 2.0, RGBColor(0x0F, 0x3A, 0x55))
add_text_box(slide, "🎵 音ってふしぎ！", 0.5, 1.0, 12.3, 1.2,
             font_size=40, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
add_text_box(slide, "世界に１つだけの楽器を作ろう！", 0.5, 2.1, 12.3, 1.0,
             font_size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, "〜 糸電話・気柱共鳴・カリンバで音のひみつをさぐれ！ 〜", 0.5, 3.1, 12.3, 0.7,
             font_size=20, color=RGBColor(0xAA,0xCC,0xFF), align=PP_ALIGN.CENTER)
add_text_box(slide, "理科教育サークル　夏の自由研究", 0.5, 6.2, 12.3, 0.7,
             font_size=18, color=C_WHITE, align=PP_ALIGN.CENTER)

# ─── SLIDE 2: もくじ ─────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, C_LIGHT_BG)
add_rect(slide, 0, 0, 13.33, 1.3, C_TITLE_BG)
add_text_box(slide, "📋  もくじ", 0.4, 0.25, 12, 0.8,
             font_size=32, bold=True, color=C_WHITE)
# 4 section cards
sections = [
    (C_PART1_BG,   "① 導入・糸電話", "音ってなに？振動を体験しよう！", "〜10分〜"),
    (C_PART2_BG,   "② 気柱共鳴",     "音は波！大きさと高さのひみつ",  "〜10分〜"),
    (C_PART3_BG,   "③ カリンバ制作", "自分だけの楽器を作ろう！",       "〜30分〜"),
    (C_SUMMARY_BG, "④ まとめ",        "わかったこと・感想",              "〜5分〜"),
]
for i, (col, num, desc, time) in enumerate(sections):
    x = 0.35 + i * 3.2
    add_rect(slide, x, 1.5, 3.0, 5.5, col)
    add_text_box(slide, num,  x+0.1, 1.6,  2.8, 0.8,  font_size=18, bold=True, color=C_WHITE)
    add_text_box(slide, desc, x+0.1, 2.4,  2.8, 2.5,  font_size=15, color=C_WHITE)
    add_text_box(slide, time, x+0.1, 6.5,  2.8, 0.4,  font_size=14, color=C_ACCENT, bold=True)

# ─── SLIDE 3: 今日のストーリー ───────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, C_LIGHT_BG)
add_rect(slide, 0, 0, 13.33, 1.3, C_TITLE_BG)
add_text_box(slide, "🗺  今日のストーリー", 0.4, 0.25, 12, 0.8, font_size=32, bold=True, color=C_WHITE)
add_text_box(slide, "音のひみつをさぐれ！３つのミッション",
             0.5, 1.4, 12.3, 0.7, font_size=24, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
arrows = [
    (C_PART1_BG,   "STEP①\n糸電話",   "音は「振動」\nで伝わる！"),
    (C_PART2_BG,   "STEP②\n気柱共鳴", "音は「波」\n長さで変わる！"),
    (C_PART3_BG,   "STEP③\nカリンバ", "自分で\n音を作る！"),
]
for i, (col, step, hint) in enumerate(arrows):
    x = 0.6 + i * 4.1
    add_rect(slide, x, 2.2, 3.5, 2.5, col)
    add_text_box(slide, step, x+0.1, 2.3, 3.3, 1.0, font_size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, hint, x+0.1, 3.3, 3.3, 1.3, font_size=18, color=C_WHITE, align=PP_ALIGN.CENTER)
    if i < 2:
        add_text_box(slide, "→", x+3.55, 3.1, 0.5, 0.8, font_size=36, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
add_text_box(slide, "すべてつながって　「音のひみつ」が完成！🎵",
             0.5, 5.2, 12.3, 0.8, font_size=22, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)

# ─── SLIDE 4: 導入質問 ───────────────────────────────────────────────
def s4(slide):
    add_text_box(slide, "みんなに質問です！", 0.5, 1.7, 12.3, 0.7, font_size=26, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    questions = [
        "🔊  「音」ってどこから来ているんだろう？",
        "🎺  楽器はなぜ音が出るの？",
        "📻  音はどうやって耳まで届くの？",
    ]
    bullet_box(slide, questions, top=2.5, font_size=24, prefix="")
    add_text_box(slide, "今日の実験で全部わかるようになるよ！🎉",
                 0.5, 6.0, 12.3, 0.8, font_size=22, bold=True, color=C_PART3_BG, align=PP_ALIGN.CENTER)
add_slide(prs, C_TITLE_BG, "導入", "音ってなに？　考えてみよう！", "❓", s4)

# ═══════════════════════════════════════════════════════════════════
# PART 1 : 糸電話 (Slides 5-14)
# ═══════════════════════════════════════════════════════════════════
P1 = "Part① 糸電話"
P1C = C_PART1_BG

def s5(slide):
    highlight_box(slide, "「糸電話」で音のひみつを調べよう！", top=1.7, bg=RGBColor(0xFF,0xE5,0xB4), font_size=26)
    add_text_box(slide, "目標：音がどのように伝わるかを体験する", 0.5, 3.4, 12.3, 0.6, font_size=22, color=C_DARK)
    add_text_box(slide, "📌 実験中に注目してほしいこと",0.5, 4.1, 12.3, 0.5, font_size=20, bold=True, color=C_TITLE_BG)
    bullet_box(slide, ["糸をつまむと声はどう変わる？", "糸がたるんでいると？", "糸なしだと届く？"], top=4.6, font_size=20)
add_slide(prs, P1C, P1, "糸電話をやってみよう！", "📞", s5)

def s6(slide):
    two_col(slide,
        ["紙コップ　2個", "たこ糸（約3〜5m）", "つまようじ"],
        ["はさみ", "セロハンテープ", "(あれば) 毛糸・ビニール紐"],
        left_head="【用意するもの】", right_head="【比べてみよう】", top=1.7)
add_slide(prs, P1C, P1, "用意するもの", "🛒", s6)

def s7(slide):
    step_box(slide, [
        "紙コップの底に　つまようじで　小さな穴をあける",
        "糸を穴に通して　内側でかた結びをする",
        "もう一方のコップも同じようにする",
        "糸をピンと張って　話しかけてみよう！",
    ], top=1.7)
add_slide(prs, P1C, P1, "作り方", "✂️", s7)

def s8(slide):
    add_text_box(slide, "やってみよう！", 0.5, 1.7, 12.3, 0.7, font_size=28, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    checks = [
        "①　ふつうに使う　→　聞こえる？",
        "②　糸をつまんでみる　→　どうなった？",
        "③　糸をたるませてみる　→　どうなった？",
        "④　糸なしで呼んでみる　→　届く？",
    ]
    bullet_box(slide, checks, top=2.5, font_size=22, prefix="")
    add_text_box(slide, "📝 結果をメモしておこう！", 0.5, 6.1, 12.3, 0.6, font_size=20, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
add_slide(prs, P1C, P1, "実験① ：やってみよう！（観察）", "🔬", s8)

def s9(slide):
    add_text_box(slide, "なぜ聞こえるの？　考えてみよう", 0.5, 1.7, 12.3, 0.7, font_size=26, bold=True, color=C_DARK)
    qs = ["糸をつまんだら聞こえにくくなったのはなぜ？",
          "糸がたるむと聞こえなくなったのはなぜ？",
          "糸がないと届かないのはなぜ？"]
    bullet_box(slide, qs, top=2.5, font_size=22)
    add_text_box(slide, "ヒント：コップに声を出した時、コップは「どう動いている」かな？",
                 0.5, 5.8, 12.3, 0.7, font_size=20, color=C_TITLE_BG, bold=True, align=PP_ALIGN.CENTER)
add_slide(prs, P1C, P1, "考察してみよう！", "💭", s9)

def s10(slide):
    highlight_box(slide, "音は「振動（しんどう）」だ！", top=1.8,
                  bg=RGBColor(0xFF,0xE0,0x80), font_size=32)
    add_text_box(slide, "振動 ＝ もの（空気・糸・コップ）が　細かく「ふるえる」こと",
                 0.5, 3.1, 12.3, 0.7, font_size=22, color=C_DARK, align=PP_ALIGN.CENTER)
    facts = ["声を出すと　のどが振動する",
             "振動が糸を通して　相手のコップに届く",
             "コップが振動して　耳に音として聞こえる"]
    bullet_box(slide, facts, top=4.0, font_size=22)
add_slide(prs, P1C, P1, "答え：音の正体", "💡", s10)

def s11(slide):
    add_text_box(slide, "振動の伝わり方イメージ", 0.5, 1.7, 12.3, 0.6, font_size=24, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    chain = ["自分の\nのど\n（振動）", "紙コップ\n①\n（振動）", "糸\n（振動\n伝わる）", "紙コップ\n②\n（振動）", "相手の\n耳\n（聞こえる）"]
    cols = [RGBColor(0xFF,0x8C,0x00), RGBColor(0xFF,0xA5,0x30), RGBColor(0xFF,0xC0,0x60),
            RGBColor(0xFF,0xA5,0x30), RGBColor(0xFF,0x8C,0x00)]
    for i, (text, col) in enumerate(zip(chain, cols)):
        x = 0.4 + i * 2.5
        add_rect(slide, x, 2.4, 2.0, 2.5, col)
        add_text_box(slide, text, x+0.05, 2.5, 1.9, 2.3, font_size=17, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        if i < 4:
            add_text_box(slide, "→", x+2.05, 3.3, 0.4, 0.7, font_size=28, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide, "糸をつまむ　→　振動がストップ　→　聞こえなくなる！",
                 0.5, 5.3, 12.3, 0.7, font_size=22, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
add_slide(prs, P1C, P1, "振動の伝わり方", "〰️", s11)

def s12(slide):
    add_text_box(slide, "空気中でも振動は伝わる！", 0.5, 1.7, 12.3, 0.6, font_size=26, bold=True, color=C_DARK)
    facts = ["声 → 空気が振動 → 耳に届く",
             "でも、遠くなるほど振動は弱くなる（音が小さくなる）",
             "糸電話では糸が振動を効率よく伝えてくれる！"]
    bullet_box(slide, facts, top=2.5, font_size=22)
    add_text_box(slide, "🌙  宇宙では音は聞こえない！（空気がないから振動しない）",
                 0.5, 5.5, 12.3, 0.7, font_size=20, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
add_slide(prs, P1C, P1, "空気でも振動は伝わる", "💨", s12)

def s13(slide):
    add_text_box(slide, "でも…", 0.5, 1.7, 12.3, 0.7, font_size=30, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide, "音の「高さ」はどこで決まるの？", 0.5, 2.5, 12.3, 0.8, font_size=34, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide, "ドレミファソラシド… 高い音・低い音、何が違うの？",
                 0.5, 3.6, 12.3, 0.7, font_size=24, color=C_DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide, "次の実験でさぐろう！ →",
                 0.5, 5.5, 12.3, 0.7, font_size=24, bold=True, color=C_PART2_BG, align=PP_ALIGN.CENTER)
add_slide(prs, P1C, P1, "まとめ① ：音は振動！　次の疑問へ…", "🔗", s13)

# ═══════════════════════════════════════════════════════════════════
# PART 2 : 気柱共鳴 (Slides 14-23)
# ═══════════════════════════════════════════════════════════════════
P2 = "Part② 気柱共鳴"
P2C = C_PART2_BG

def s14(slide):
    highlight_box(slide, "音叉＋試験管　で　気柱共鳴を観察しよう！", top=1.8, bg=RGBColor(0xD0,0xEA,0xFF), font_size=26)
    add_text_box(slide, "目標：音の高さが何で決まるかを「目で見て」理解する", 0.5, 3.2, 12.3, 0.6, font_size=22, color=C_DARK)
    add_text_box(slide, "📌 注目ポイント", 0.5, 4.0, 12.3, 0.5, font_size=20, bold=True, color=C_TITLE_BG)
    bullet_box(slide, ["音叉をどこに近づけると音が大きくなる？",
                       "試験管の長さを変えると音は変わる？"], top=4.6, font_size=20)
add_slide(prs, P2C, P2, "気柱共鳴を観察しよう！", "🎵", s14)

def s15(slide):
    two_col(slide,
        ["音叉（おんさ）　数種類", "試験管（長さが違うもの）", "水（試験管に入れる用）", "ゴムハンマー"],
        ["聴診器（あれば）", "スピーカー（演示用）", "記録用紙"],
        left_head="【用意するもの】", right_head="【あると便利】", top=1.7)
add_slide(prs, P2C, P2, "用意するもの", "🛒", s15)

def s16(slide):
    step_box(slide, [
        "音叉をゴムハンマーで叩いて音を出す",
        "試験管の口に音叉を近づける",
        "試験管の長さを変えながら　音の変化を観察する",
        "水を入れて試験管の長さを調整してみよう",
        "一番音が大きくなる長さを記録しよう！",
    ], top=1.7)
add_slide(prs, P2C, P2, "実験の方法", "📋", s16)

def s17(slide):
    add_text_box(slide, "観察してみよう！", 0.5, 1.7, 12.3, 0.6, font_size=28, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    obs = ["試験管が短い　→　音は？", "試験管が長い　→　音は？", "一番大きくなったのはどの長さのとき？"]
    bullet_box(slide, obs, top=2.5, font_size=24, prefix="🔍")
    add_text_box(slide, "📝 気づいたことをメモしよう！", 0.5, 5.8, 12.3, 0.6, font_size=20, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
add_slide(prs, P2C, P2, "演示① ：観察タイム", "👀", s17)

def s18(slide):
    highlight_box(slide, "「共鳴（きょうめい）」って何？", top=1.8, bg=RGBColor(0xD0,0xEA,0xFF), font_size=30)
    add_text_box(slide,
        "ブランコのたとえ\n→ タイミングが合うと　どんどん大きく揺れる！\n→ 音も同じ！　振動のリズムが合うと　大きく響く！",
        0.5, 3.1, 12.3, 2.5, font_size=22, color=C_DARK)
    add_text_box(slide, "試験管の空気の長さが　音叉の振動に「ぴったり」合うと　大きく鳴る！",
                 0.5, 5.8, 12.3, 0.7, font_size=20, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
add_slide(prs, P2C, P2, "共鳴とは？", "💡", s18)

def s19(slide):
    add_text_box(slide, "音は「波（なみ）」だ！", 0.5, 1.7, 12.3, 0.7, font_size=30, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide,
        "振動が広がっていく様子 ＝ 「波」\n\n"
        "　波の長さのことを「波長（はちょう）」という\n"
        "　1秒間に何回振動するか ＝ 「振動数（しんどうすう）」（単位：Hz ヘルツ）",
        0.5, 2.5, 12.3, 3.0, font_size=22, color=C_DARK)
    add_text_box(slide, "振動数が多い　→　高い音　　　　振動数が少ない　→　低い音",
                 0.5, 5.7, 12.3, 0.7, font_size=22, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
add_slide(prs, P2C, P2, "音の正体：「波」と「振動数」", "〰️", s19)

def s20(slide):
    add_text_box(slide, "試験管の長さと音の高さの関係", 0.5, 1.7, 12.3, 0.6, font_size=26, bold=True, color=C_DARK)
    rows = [
        ("長い（空気の柱が長い）", "波長が長い", "振動数が少ない", "低い音"),
        ("短い（空気の柱が短い）", "波長が短い", "振動数が多い", "高い音"),
    ]
    headers = ["試験管の長さ", "波長", "振動数", "音の高さ"]
    hx = [0.4, 3.5, 6.5, 9.8]
    add_rect(slide, 0.3, 2.4, 12.7, 0.55, C_PART2_BG)
    for col, h in zip(hx, headers):
        add_text_box(slide, h, col, 2.45, 3.0, 0.45, font_size=18, bold=True, color=C_WHITE)
    for ri, row in enumerate(rows):
        bg = RGBColor(0xE8,0xF5,0xFF) if ri%2==0 else C_WHITE
        add_rect(slide, 0.3, 3.0+ri*0.8, 12.7, 0.75, bg)
        for col, cell in zip(hx, row):
            add_text_box(slide, cell, col, 3.05+ri*0.8, 3.0, 0.65, font_size=18, color=C_DARK)
    add_text_box(slide, "楽器でも全く同じ！　管の長さで音の高さが変わる！",
                 0.5, 5.5, 12.3, 0.7, font_size=22, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
add_slide(prs, P2C, P2, "波長と音の高さの関係", "📊", s20)

def s21(slide):
    add_text_box(slide, "身近な楽器との共通点", 0.5, 1.7, 12.3, 0.6, font_size=26, bold=True, color=C_DARK)
    inst = [
        "🎺  トランペット　→　管の長さで音が変わる",
        "🎸  ギター　→　弦の長さで音が変わる",
        "🪗  リコーダー　→　穴をふさいで空気の長さを変える",
        "🥁  太鼓　→　皮の大きさで音が変わる",
    ]
    bullet_box(slide, inst, top=2.5, font_size=22, prefix="")
    add_text_box(slide, "全部「振動する部分の長さ・大きさ」で音が決まっている！",
                 0.5, 6.0, 12.3, 0.7, font_size=22, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
add_slide(prs, P2C, P2, "楽器も同じしくみ！", "🎶", s21)

def s22(slide):
    add_text_box(slide, "まとめ②", 0.5, 1.7, 12.3, 0.6, font_size=22, bold=True, color=C_DARK)
    highlight_box(slide, "音の高さは「振動の数（振動数）」で決まる！", top=2.5, bg=RGBColor(0xD0,0xEA,0xFF), font_size=28)
    pts = ["空気の柱が長い → 振動数が少ない → 低い音",
           "空気の柱が短い → 振動数が多い → 高い音",
           "これを使えば　自分で音を作れる！"]
    bullet_box(slide, pts, top=3.8, font_size=22)
    add_text_box(slide, "→ さあ、カリンバを作ろう！🎵",
                 0.5, 6.1, 12.3, 0.7, font_size=24, bold=True, color=C_PART3_BG, align=PP_ALIGN.CENTER)
add_slide(prs, P2C, P2, "まとめ② ：音の高さのひみつ", "✅", s22)

# ═══════════════════════════════════════════════════════════════════
# PART 3 : カリンバ制作 (Slides 23-36)
# ═══════════════════════════════════════════════════════════════════
P3 = "Part③ カリンバ制作"
P3C = C_PART3_BG

def s23(slide):
    highlight_box(slide, "いよいよ！　世界に１つだけの楽器を作ろう！", top=1.8, bg=RGBColor(0xD0,0xFF,0xD8), font_size=28)
    add_text_box(slide,
        "今まで学んだこと　→　全部カリンバに入っている！\n\n"
        "　①　音は「振動」で伝わる　（糸電話）\n"
        "　②　振動の数で音の高さが変わる　（気柱共鳴）\n"
        "　③　自分でキーの長さを変えて　自分だけの音を作る！",
        0.5, 3.0, 12.3, 3.5, font_size=22, color=C_DARK)
add_slide(prs, P3C, P3, "いよいよ楽器制作！", "🎹", s23)

def s24(slide):
    add_text_box(slide, "カリンバってなに？", 0.5, 1.7, 12.3, 0.6, font_size=28, bold=True, color=C_DARK)
    add_text_box(slide,
        "アフリカ生まれの民族楽器\n"
        "英語では　「サムピアノ」（親指ピアノ）ともよばれる\n\n"
        "金属のキー（板）を親指ではじいて音を出す\n"
        "キーの長さを変えることで　音の高さが決まる！",
        0.5, 2.5, 12.3, 3.0, font_size=22, color=C_DARK)
    add_text_box(slide, "→ まさに「気柱共鳴」で学んだ　「長さ→音の高さ」の原理！",
                 0.5, 5.9, 12.3, 0.7, font_size=20, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
add_slide(prs, P3C, P3, "カリンバってなに？", "🌍", s24)

def s25(slide):
    two_col(slide,
        ["木の板（台座）", "金属のキー　8〜10本", "固定用の金属棒（ブリッジ）"],
        ["ボルト・ナット（固定用）", "やすり（あれば）", "デコレーション用シール"],
        left_head="【作るもの】", right_head="【あると便利】", top=1.7)
    add_text_box(slide, "📦 材料はあらかじめ用意してあります！",
                 0.5, 6.1, 12.3, 0.6, font_size=20, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
add_slide(prs, P3C, P3, "用意するもの", "🛒", s25)

def s26(slide):
    step_box(slide, [
        "木の板に　ブリッジ（金属棒）を置く場所を確認する",
        "金属キーを　長いもの → 短いもの　の順にならべる",
        "ブリッジをのせて　ボルトでしっかり固定する",
        "各キーの長さを少しずつ調整する（ここが大事！）",
    ], top=1.7)
    add_text_box(slide, "ブリッジの外側に出る長さ　＝　音の高さ　を決める！",
                 0.5, 6.1, 12.3, 0.6, font_size=20, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
add_slide(prs, P3C, P3, "作り方①：土台を作ろう", "🔨", s26)

def s27(slide):
    step_box(slide, [
        "キーを親指で軽くはじいてみよう",
        "音が出たかな？（ちゃんと固定できてるかな）",
        "キーを前後にずらして　音の高さを調整する",
        "ドレミ… になるように　少しずつ動かそう！",
    ], top=1.7)
    add_text_box(slide, "キーを長くする → 低い音　　キーを短くする → 高い音",
                 0.5, 6.1, 12.3, 0.6, font_size=20, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
add_slide(prs, P3C, P3, "作り方②：音を出して調整しよう", "🎵", s27)

def s28(slide):
    add_text_box(slide, "なぜ長さで音が変わるの？", 0.5, 1.7, 12.3, 0.6, font_size=26, bold=True, color=C_DARK)
    add_text_box(slide,
        "キーが長い\n→ はじかれた時の「振動」がゆっくり（振動数が少ない）\n→ 低い音\n\n"
        "キーが短い\n→ 振動が速い（振動数が多い）\n→ 高い音",
        0.5, 2.5, 12.3, 3.3, font_size=22, color=C_DARK)
    add_text_box(slide, "気柱共鳴と　まったく同じ原理！",
                 0.5, 6.0, 12.3, 0.7, font_size=22, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
add_slide(prs, P3C, P3, "キーの長さと音の高さ", "🔗", s28)

def s29(slide):
    add_text_box(slide, "糸電話・気柱共鳴・カリンバ　全部つながっている！", 0.5, 1.7, 12.3, 0.6, font_size=22, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    links = [
        (C_PART1_BG, "糸電話", "音は振動で\n伝わる"),
        (C_PART2_BG, "気柱共鳴", "振動数で\n音の高さが決まる"),
        (C_PART3_BG, "カリンバ", "キーの長さで\n振動数を調整"),
    ]
    for i, (col, label, desc) in enumerate(links):
        x = 0.8 + i * 4.0
        add_rect(slide, x, 2.5, 3.5, 3.5, col)
        add_text_box(slide, label, x+0.1, 2.6, 3.3, 0.8, font_size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, desc,  x+0.1, 3.4, 3.3, 2.5, font_size=20, color=C_WHITE, align=PP_ALIGN.CENTER)
        if i < 2:
            add_text_box(slide, "→", x+3.55, 3.7, 0.4, 0.8, font_size=36, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
add_slide(prs, P3C, P3, "３つの実験のつながり", "🔗", s29)

def s30(slide):
    add_text_box(slide, "デコレーションしよう！", 0.5, 1.7, 12.3, 0.6, font_size=28, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide,
        "世界に１つだけの　自分だけのカリンバを飾ろう！\n\n"
        "　・シールを貼る\n"
        "　・名前を書く\n"
        "　・音の名前（ドレミ）を書いてみよう",
        0.5, 2.5, 12.3, 3.0, font_size=22, color=C_DARK)
    add_text_box(slide, "完成したら　弾いてみよう！ 🎵",
                 0.5, 6.0, 12.3, 0.7, font_size=24, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
add_slide(prs, P3C, P3, "自分だけの楽器を完成させよう！", "🎨", s30)

def s31(slide):
    add_text_box(slide, "みんなで演奏してみよう！🎵", 0.5, 1.7, 12.3, 0.7, font_size=30, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide,
        "「きらきら星」の弾き方\n\n"
        "ド　ド　ソ　ソ　ラ　ラ　ソー\n"
        "ファ　ファ　ミ　ミ　レ　レ　ドー",
        0.5, 2.6, 12.3, 2.8, font_size=24, color=C_DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide, "上手に弾けたかな？ どんな音がしたか教えてね！",
                 0.5, 5.8, 12.3, 0.7, font_size=22, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
add_slide(prs, P3C, P3, "みんなで演奏！", "🎶", s31)

# ═══════════════════════════════════════════════════════════════════
# まとめ (Slides 32-38)
# ═══════════════════════════════════════════════════════════════════
SC = C_SUMMARY_BG
SL = "まとめ"

def s32(slide):
    add_text_box(slide, "今日わかったこと", 0.5, 1.7, 12.3, 0.6, font_size=28, bold=True, color=C_DARK)
    pts = [
        "① 音は「振動（しんどう）」で伝わる",
        "② 振動の数（振動数）が多いほど　高い音になる",
        "③ 振動する部分の「長さ」で　音の高さが変わる",
        "④ これをうまく使えば　楽器が作れる！",
    ]
    bullet_box(slide, pts, top=2.5, font_size=24)
add_slide(prs, SC, SL, "今日わかったこと", "📚", s32)

def s33(slide):
    add_text_box(slide, "音のひみつ まとめ", 0.5, 1.7, 12.3, 0.6, font_size=28, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    highlight_box(slide,
        "音 ＝ 振動　　　高い音 ＝ 速い振動　　　低い音 ＝ 遅い振動",
        top=2.5, bg=RGBColor(0xEE,0xDD,0xFF), font_size=24)
    add_text_box(slide,
        "身の周りの楽器　→　全部「振動の長さ・速さ」で音が決まっている\n"
        "リコーダー・ギター・ピアノ・太鼓　全部同じ原理！",
        0.5, 4.2, 12.3, 2.0, font_size=22, color=C_DARK)
add_slide(prs, SC, SL, "音のひみつ まとめ", "🔑", s33)

def s34(slide):
    add_text_box(slide, "自由研究にまとめてみよう！", 0.5, 1.7, 12.3, 0.6, font_size=26, bold=True, color=C_DARK)
    pts = [
        "①  テーマ：「音の伝わり方と楽器のしくみ」",
        "②  仮説：音は○○で伝わる、と思う",
        "③  実験方法：糸電話・気柱共鳴・カリンバ",
        "④  結果：わかったこと",
        "⑤  考察：なぜそうなった？",
        "⑥  感想：不思議に思ったこと・もっと調べたいこと",
    ]
    bullet_box(slide, pts, top=2.5, font_size=20)
add_slide(prs, SC, SL, "自由研究にまとめてみよう！", "📝", s34)

def s35(slide):
    add_text_box(slide, "感想を書こう！", 0.5, 1.7, 12.3, 0.6, font_size=28, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    qs = [
        "今日一番おもしろかった実験はどれ？",
        "音についてどんなことがわかった？",
        "もっと調べてみたいことはある？",
        "カリンバを弾いてみてどうだった？",
    ]
    bullet_box(slide, qs, top=2.5, font_size=22)
add_slide(prs, SC, SL, "感想・気づきを教えてね！", "💬", s35)

def s36(slide):
    add_text_box(slide, "もっと知りたい人へ！", 0.5, 1.7, 12.3, 0.6, font_size=28, bold=True, color=C_DARK)
    more = [
        "🔬  音の速さを測ってみよう（音速：約340m/秒）",
        "🎼  ピアノの弦を観察してみよう",
        "🌊  水を入れたグラスを指でこすると…？",
        "📱  スマホの音声波形を見てみよう（無料アプリあり）",
    ]
    bullet_box(slide, more, top=2.5, font_size=22, prefix="")
    add_text_box(slide, "音の世界はまだまだ深い！どんどん探求しよう！",
                 0.5, 6.0, 12.3, 0.7, font_size=20, bold=True, color=C_TITLE_BG, align=PP_ALIGN.CENTER)
add_slide(prs, SC, SL, "もっと探求しよう！", "🚀", s36)

# ─── SLIDE 37: おわりに ───────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, C_TITLE_BG)
add_rect(slide, 0, 5.0, 13.33, 2.5, RGBColor(0x0F, 0x3A, 0x55))
add_text_box(slide, "🎵 今日はありがとう！", 0.5, 1.0, 12.3, 1.0,
             font_size=40, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
add_text_box(slide, "音のひみつ　わかったかな？", 0.5, 2.1, 12.3, 0.8,
             font_size=32, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, "作ったカリンバ　大切にしてね！\n家でも友だちに教えてあげてね！🎹",
             0.5, 3.0, 12.3, 1.5, font_size=24, color=RGBColor(0xAA,0xCC,0xFF), align=PP_ALIGN.CENTER)
add_text_box(slide, "理科教育サークル", 0.5, 6.2, 12.3, 0.6,
             font_size=18, color=C_WHITE, align=PP_ALIGN.CENTER)

# ─── Save ────────────────────────────────────────────────────────────
out_path = "/home/user/nust-room-search/音ってふしぎ_スライド概要.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")

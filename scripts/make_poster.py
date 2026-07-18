#!/usr/bin/env python3
"""RoomRadar 学内掲示ポスター生成スクリプト。

正式版（--variant official）と静か版（--variant shizuka）を生成できる。
デザイン: 濃紺ヘッダー/フッター帯・実機スクショのスマホモック（自然な3D・
影は下方向のみ）・Instagram QRを主役に配置・暖色→寒色のブランドグラデ背景。

使い方:
  python scripts/make_poster.py                       # 正式版 → assets/posters/roomradar-poster-pro.png
  python scripts/make_poster.py --print-files         # PNGに加えPDF/PPTXも出力
  python scripts/make_poster.py --variant shizuka --out assets/posters/roomradar-poster-shizuka.png

必要パッケージ: qrcode, pillow, playwright, python-pptx（--print-files時）
素材: rr_logo.png / assets/posters/handoff/{app-phone-screenshot,qr-instagram-branded}.png

注意: 静か版はネタ文言を含む仮バージョン。掲示先での運用判断のため、成果物は
assets/posters/roomradar-poster-shizuka.{png,pdf,pptx} としてリポジトリに保存している。
"""
import argparse
import base64
import io
import pathlib

import qrcode
from qrcode.constants import ERROR_CORRECT_H
from playwright.sync_api import sync_playwright

ROOT = pathlib.Path(__file__).resolve().parents[1]
LP_URL = "https://nu-roomradar.github.io/nust-room-search/"
CHROME = "/opt/pw-browsers/chromium-1194/chrome-linux/chrome"

LOGO_PATH = ROOT / "rr_logo.png"
SHOT_PATH = ROOT / "assets/posters/handoff/app-phone-screenshot.png"
IG_QR_PATH = ROOT / "assets/posters/handoff/qr-instagram-branded.png"
DEFAULT_OUT = ROOT / "assets/posters/roomradar-poster-pro.png"

# バリアントごとの文言（eyebrow・説明文）
VARIANTS = {
    "official": {
        "eyebrow": "次のコマ、どこで過ごす？",
        "lead": (
            "曜日・時限・校舎を選ぶだけ。時間割データをもとに、授業が入っていない教室を"
            "リアルタイムで表示します。もう教室を探して歩き回らない。"
            "静かな教室で、集中して勉強・作業ができます。"
        ),
    },
    "shizuka": {
        "eyebrow": "教室では“静か”にしろ",
        "lead": (
            "曜日・時限・校舎を選ぶだけ。時間割データをもとに、授業が入っていない教室を"
            "リアルタイムで表示します。もう教室を探して歩き回らない。"
            '<span class="kw">“静か”</span>な教室で、集中して勉強・作業ができます。'
        ),
    },
}


def datauri(path):
    return "data:image/png;base64," + base64.b64encode(pathlib.Path(path).read_bytes()).decode()


def qr_datauri(data):
    qr = qrcode.QRCode(error_correction=ERROR_CORRECT_H, box_size=20, border=2)
    qr.add_data(data)
    qr.make(fit=True)
    img = qr.make_image(fill_color="#101828", back_color="white").convert("RGB")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return "data:image/png;base64," + base64.b64encode(buf.getvalue()).decode()


IC_SEARCH = '<svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="1.7" stroke-linecap="round" stroke-linejoin="round"><circle cx="11" cy="11" r="7"/><line x1="21" y1="21" x2="16.65" y2="16.65"/></svg>'
IC_FREE = '<svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="1.7" stroke-linecap="round" stroke-linejoin="round"><path d="M20.59 13.41l-7.17 7.17a2 2 0 0 1-2.83 0L2 12V2h10l8.59 8.59a2 2 0 0 1 0 2.82z"/><line x1="7" y1="7" x2="7.01" y2="7"/></svg>'
IC_BUILD = '<svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="1.7" stroke-linecap="round" stroke-linejoin="round"><line x1="3" y1="21" x2="21" y2="21"/><path d="M4 21V8l8-5 8 5v13"/><line x1="9" y1="21" x2="9" y2="13"/><line x1="15" y1="21" x2="15" y2="13"/><line x1="4" y1="9.5" x2="20" y2="9.5"/></svg>'
IC_IG = '<svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="1.7" stroke-linecap="round" stroke-linejoin="round"><rect x="3" y="3" width="18" height="18" rx="5"/><circle cx="12" cy="12" r="4"/><line x1="17.5" y1="6.5" x2="17.51" y2="6.5"/></svg>'
IC_APP = '<svg viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="1.7" stroke-linecap="round" stroke-linejoin="round"><rect x="5" y="2" width="14" height="20" rx="3"/><line x1="12" y1="18" x2="12" y2="18"/></svg>'


def build_html(variant):
    v = VARIANTS[variant]
    logo = datauri(LOGO_PATH)
    shot = datauri(SHOT_PATH)
    qr_lp = qr_datauri(LP_URL)
    qr_ig = datauri(IG_QR_PATH)
    return f"""<!DOCTYPE html><html lang="ja"><head><meta charset="UTF-8">
<link href="https://fonts.googleapis.com/css2?family=Syne:wght@700;800&family=Noto+Sans+JP:wght@400;500;700;900&display=swap" rel="stylesheet">
<style>
  * {{ margin:0; padding:0; box-sizing:border-box; }}
  :root {{ --ink:#101828; --muted:#6b7280; --accent:#2563eb; --accent2:#f97316;
    --navy:#14203a; --navy2:#1e2f52; --border:rgba(16,24,40,0.1); }}
  html,body {{ width:1080px; height:1528px; }}
  body {{ font-family:'Noto Sans JP',sans-serif; color:var(--ink);
    background:
      radial-gradient(1250px 780px at 86% 12%, rgba(249,115,22,0.24), transparent 62%),
      radial-gradient(1150px 800px at 8% 94%, rgba(37,99,235,0.24), transparent 62%),
      linear-gradient(160deg, #ffdfc4 0%, #f0e9f2 50%, #cddcfa 100%);
    display:flex; flex-direction:column; }}

  /* ── ヘッダー帯 ── */
  .top {{ position:relative; overflow:hidden; height:184px; flex-shrink:0;
    background:linear-gradient(120deg,var(--navy) 0%,var(--navy2) 100%);
    display:flex; align-items:center; padding:0 72px; }}
  .top .radar {{ position:absolute; right:-90px; top:-140px; width:440px; height:440px; opacity:0.16; }}
  .top .radar span {{ position:absolute; border:2px solid #7fa8ff; border-radius:50%; inset:0; }}
  .top .radar span:nth-child(2){{ inset:60px; }} .top .radar span:nth-child(3){{ inset:120px; }} .top .radar span:nth-child(4){{ inset:180px; }}
  .top img {{ height:92px; margin-right:22px; filter:drop-shadow(0 4px 10px rgba(0,0,0,0.3)); }}
  .brandtxt .wm {{ font-family:'Syne',sans-serif; font-weight:800; font-size:2.7rem; color:#fff; letter-spacing:-0.02em; line-height:1; }}
  .brandtxt .wm .b {{ color:#7fa8ff; }}
  .brandtxt .tg {{ color:#b9c4dc; font-size:1.02rem; font-weight:500; margin-top:8px; letter-spacing:0.05em; }}

  /* ── 中央（均等配分） ── */
  .main {{ flex:1; display:flex; flex-direction:column; justify-content:space-evenly; padding:8px 0; }}

  /* ── ヒーロー ── */
  .hero {{ display:flex; padding:0 72px; gap:26px; align-items:center; position:relative; z-index:1; }}
  .hero-l {{ flex:1.1; }}
  .eyebrow {{ display:inline-block; font-weight:700; font-size:1rem; color:var(--accent);
    background:rgba(37,99,235,0.09); border:1px solid rgba(37,99,235,0.22); padding:8px 18px; border-radius:99px; letter-spacing:0.04em; }}
  h1 {{ font-weight:900; font-size:4.1rem; line-height:1.24; margin-top:26px; letter-spacing:-0.02em; }}
  h1 .hl {{ color:var(--accent); }}
  .lead {{ font-size:1.34rem; color:#4b5563; line-height:1.75; margin-top:26px; font-weight:500; max-width:530px; }}
  .lead .kw {{ white-space:nowrap; font-weight:800; color:var(--accent); }}
  .hero-r {{ flex:0.9; display:flex; justify-content:center; align-items:center; position:relative; perspective:1700px; }}

  /* ── 接地シャドウ ── */
  .phone-shadow {{ position:absolute; bottom:0; left:50%; transform:translateX(-50%); z-index:0;
    width:240px; height:40px; border-radius:50%;
    background:radial-gradient(ellipse at center, rgba(16,24,40,0.35) 0%, rgba(16,24,40,0.13) 48%, rgba(16,24,40,0) 72%);
    filter:blur(9px); }}

  /* ── スマホモック（自然な正面寄り3D・影は下方向のみ） ── */
  .phone {{ position:relative; z-index:1; width:296px; height:576px; border-radius:44px; background:#05070d;
    padding:13px;
    box-shadow:
      0 30px 48px rgba(16,24,40,0.28),
      0 12px 22px rgba(16,24,40,0.16),
      inset 0 0 0 1px rgba(255,255,255,0.07),
      inset 2px 3px 6px rgba(255,255,255,0.10);
    transform:rotateY(-5deg) rotateX(2.5deg) rotate(0.6deg);
    transform-style:preserve-3d;
    border:1px solid rgba(255,255,255,0.12); }}
  .phone .screen {{ position:relative; width:100%; height:100%; border-radius:32px; overflow:hidden; background:#0a0a0f; }}
  .phone .screen img {{ width:100%; display:block; object-fit:cover; object-position:top; }}
  .phone .glare {{ position:absolute; inset:13px; border-radius:32px; z-index:2; pointer-events:none;
    background:linear-gradient(122deg, rgba(255,255,255,0.18) 0%, rgba(255,255,255,0.06) 16%, rgba(255,255,255,0) 38%); }}
  .phone .notch {{ position:absolute; top:22px; left:50%; transform:translateX(-50%); width:120px; height:26px; background:#05070d; border-radius:0 0 16px 16px; z-index:3; }}

  /* ── 特徴3点 ── */
  .feats {{ display:flex; gap:22px; padding:0 72px; position:relative; z-index:3; }}
  .feat {{ flex:1; display:flex; gap:17px; align-items:center;
    background:#fff; border:1px solid var(--border); border-radius:20px; padding:30px 26px;
    box-shadow:0 6px 20px rgba(16,24,40,0.05); }}
  .feat .ico {{ width:56px; height:56px; flex-shrink:0; border-radius:15px; display:flex; align-items:center; justify-content:center;
    background:rgba(37,99,235,0.08); color:var(--accent); }}
  .feat .ico svg {{ width:30px; height:30px; }}
  .feat.o .ico {{ background:rgba(249,115,22,0.1); color:var(--accent2); }}
  .feat .t {{ font-weight:800; font-size:1.3rem; line-height:1.3; white-space:nowrap; }}
  .feat .d {{ font-size:0.98rem; color:var(--muted); margin-top:7px; line-height:1.5; }}

  /* ── QR（Instagramを主役に） ── */
  .qrwrap {{ padding:0 72px; }}
  .qrhead {{ text-align:center; font-weight:800; font-size:1.32rem; margin-bottom:22px; }}
  .qrhead .em {{ color:var(--accent2); }}
  .qrs {{ display:flex; gap:26px; align-items:stretch; }}
  .qr-primary {{ flex:1.35; display:flex; align-items:center; gap:28px; position:relative;
    background:#fff; border:2px solid rgba(249,115,22,0.4); border-radius:22px; padding:26px 30px;
    box-shadow:0 18px 44px rgba(249,115,22,0.16); }}
  .qr-primary .badge {{ position:absolute; top:-16px; left:30px; background:var(--accent2); color:#fff;
    font-weight:800; font-size:0.86rem; padding:6px 16px; border-radius:99px; box-shadow:0 6px 16px rgba(249,115,22,0.35); }}
  .qr-primary .qbox {{ width:224px; flex-shrink:0; }}
  .qr-primary .qbox img {{ width:100%; height:auto; display:block; }}
  .qr-primary .k {{ display:flex; align-items:center; gap:9px; font-weight:900; font-size:1.3rem; white-space:nowrap; }}
  .qr-primary .k svg {{ width:24px; height:24px; color:var(--accent2); flex-shrink:0; }}
  .qr-primary .handle {{ font-family:'Syne',sans-serif; font-weight:800; font-size:1.62rem; margin-top:8px; white-space:nowrap; }}
  .qr-primary .s {{ font-size:1rem; color:var(--muted); margin-top:12px; line-height:1.5; }}
  .qr-sec {{ flex:0.85; display:flex; flex-direction:column; align-items:center; text-align:center;
    background:#fff; border:1px solid var(--border); border-radius:22px; padding:24px 24px; }}
  .qr-sec .qbox {{ width:150px; height:150px; border-radius:14px; padding:8px; background:#fff; border:1px solid var(--border); }}
  .qr-sec .qbox img {{ width:100%; height:100%; display:block; }}
  .qr-sec .k {{ display:flex; align-items:center; gap:8px; font-weight:800; font-size:1.12rem; margin-top:16px; }}
  .qr-sec .k svg {{ width:20px; height:20px; color:var(--accent); }}
  .qr-sec .s {{ font-size:0.9rem; color:var(--muted); margin-top:7px; }}

  /* ── フッター帯 ── */
  .foot {{ background:linear-gradient(120deg,var(--navy) 0%,var(--navy2) 100%);
    padding:24px 72px 28px; color:#fff; flex-shrink:0; }}
  .foot .note {{ display:inline-flex; align-items:center; gap:11px; font-size:0.98rem; color:#dbe3f0; margin-bottom:16px;
    background:rgba(249,115,22,0.14); border:1px solid rgba(249,115,22,0.42); border-radius:10px; padding:10px 15px; }}
  .foot .note .tg {{ font-family:'Syne',sans-serif; font-weight:800; font-size:0.76rem; letter-spacing:0.12em;
    background:var(--accent2); color:#fff; border-radius:6px; padding:4px 11px; }}
  .foot .note strong {{ color:#ffd9b8; font-weight:800; }}
  .foot .row {{ display:flex; align-items:center; gap:20px; flex-wrap:wrap; }}
  .foot .org {{ font-weight:800; font-size:1.4rem; }}
  .foot .hash {{ font-weight:800; color:#7fa8ff; font-size:1.05rem; }}
  .foot .url {{ font-family:'Syne',sans-serif; font-size:0.98rem; color:#9fb0d0; margin-left:auto; }}
</style></head>
<body>
  <div class="top">
    <div class="radar"><span></span><span></span><span></span><span></span></div>
    <img src="{logo}" alt="RoomRadar">
    <div class="brandtxt">
      <div class="wm">Room<span class="b">Radar</span></div>
      <div class="tg">日本大学理工学部 · 空き教室リアルタイム検索</div>
    </div>
  </div>

  <div class="main">
    <div class="hero">
      <div class="hero-l">
        <span class="eyebrow">{v["eyebrow"]}</span>
        <h1>空いてる教室、<br><span class="hl">3秒</span>で見つかる。</h1>
        <p class="lead">{v["lead"]}</p>
      </div>
      <div class="hero-r">
        <div class="phone-shadow"></div>
        <div class="phone"><div class="notch"></div><div class="screen"><img src="{shot}" alt="RoomRadar app"></div><div class="glare"></div></div>
      </div>
    </div>

    <div class="feats">
      <div class="feat"><div class="ico">{IC_SEARCH}</div><div><div class="t">リアルタイム検索</div><div class="d">曜日・時限・校舎を選ぶだけ</div></div></div>
      <div class="feat o"><div class="ico">{IC_FREE}</div><div><div class="t">登録不要・無料</div><div class="d">ブラウザだけですぐ使える</div></div></div>
      <div class="feat"><div class="ico">{IC_BUILD}</div><div><div class="t">3校舎に対応</div><div class="d">タワースコラ / 駿河台 / 船橋</div></div></div>
    </div>

    <div class="qrwrap">
      <div class="qrhead">スマホのカメラで <span class="em">まずはInstagramをフォロー</span> ！</div>
      <div class="qrs">
        <div class="qr-primary">
          <div class="badge">まずはこちら</div>
          <div class="qbox"><img src="{qr_ig}" alt="Instagram QR"></div>
          <div>
            <div class="k">{IC_IG}Instagramをフォロー</div>
            <div class="handle">@roomradar_nust</div>
            <div class="s">最新情報・新機能・お役立ちTipsを配信中</div>
          </div>
        </div>
        <div class="qr-sec">
          <div class="qbox"><img src="{qr_lp}" alt="LP QR"></div>
          <div class="k">{IC_APP}アプリを使う</div>
          <div class="s">空き教室を今すぐ検索</div>
        </div>
      </div>
    </div>
  </div>

  <div class="foot">
    <div class="note"><span class="tg">TEST</span><span>本サービスは<strong>テスト運用中の非公式サービス</strong>です（学生が開発・運営／大学公式のものではありません）</span></div>
    <div class="row">
      <span class="org">日本大学自主創造プロジェクト</span>
      <span class="hash">#日大生プロジェクト</span>
      <span class="url">nu-roomradar.github.io/nust-room-search</span>
    </div>
  </div>
</body></html>"""


def render(html, out_png):
    tmp_html = pathlib.Path("/tmp/roomradar_poster_render.html")
    tmp_html.write_text(html, encoding="utf-8")
    with sync_playwright() as p:
        browser = p.chromium.launch(executable_path=CHROME)
        # device_scale_factor=3 → 3240x4584 ≒ A4 392dpi（印刷グレード）
        page = browser.new_page(viewport={"width": 1080, "height": 1528}, device_scale_factor=3)
        page.goto(tmp_html.as_uri(), wait_until="networkidle", timeout=30000)
        page.wait_for_timeout(1200)
        out_png.parent.mkdir(parents=True, exist_ok=True)
        page.screenshot(path=str(out_png))
        browser.close()


def export_print_files(png_path):
    """PNGと同じ場所に、A4縦のPDFとPPTXを書き出す。"""
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Mm

    img = Image.open(png_path).convert("RGB")
    dpi = img.size[0] / (210 / 25.4)  # 幅を210mmとしたときのdpi
    pdf_path = png_path.with_suffix(".pdf")
    img.save(pdf_path, "PDF", resolution=dpi)

    prs = Presentation()
    prs.slide_width = Mm(210)
    prs.slide_height = Mm(297)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(png_path), Mm(0), Mm(0), width=Mm(210), height=Mm(297))
    pptx_path = png_path.with_suffix(".pptx")
    prs.save(pptx_path)
    return pdf_path, pptx_path


def main():
    ap = argparse.ArgumentParser(description="RoomRadar 学内掲示ポスターを生成")
    ap.add_argument("--variant", choices=list(VARIANTS), default="official")
    ap.add_argument("--out", type=pathlib.Path, default=None,
                    help="出力PNGパス（省略時: 正式版はassets/posters/、静か版は/tmp/）")
    ap.add_argument("--print-files", action="store_true", help="PDF/PPTXも同じ場所に出力")
    args = ap.parse_args()

    out = args.out
    if out is None:
        out = DEFAULT_OUT if args.variant == "official" else pathlib.Path("/tmp/roomradar-poster-shizuka.png")

    render(build_html(args.variant), out)
    print("saved", out)
    if args.print_files:
        for p in export_print_files(out):
            print("saved", p)


if __name__ == "__main__":
    main()
